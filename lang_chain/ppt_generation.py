# -*- coding: utf-8 -*-
# @Time    : 2024/5/12 15:55
# @Author  : nongbin
# @FileName: ppt_generation.py
# @Software: PyCharm
# @Affiliation: tfswufe.edu.cn
import datetime
import hashlib
import os
import time
from pathlib import Path
from typing import Dict

from icecream import ic
from pptx import Presentation

from env import get_app_root
from logger import Logger
from utils.file_util import clear_files_by_timediff
from utils.schedule import get_scheduler

_logger: Logger = Logger(Path(__file__).name)

_OUTPUT_DIR = os.path.join(get_app_root(), "data/cache/ppt")

# 如果文件夹路径不存在，先创建
if not os.path.exists(_OUTPUT_DIR):
    os.makedirs(_OUTPUT_DIR)


def get_file_path(text):
    file_name = hashlib.sha256(text.encode("utf-8")).hexdigest() ## 也可以使用uuid
    return os.path.join(_OUTPUT_DIR, f"{file_name}.pptx")


def generate(ppt_content: Dict) -> str:
    ppt = Presentation()

    # PPT首页
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout
    slide.placeholders[0].text = ppt_content['title']
    slide.placeholders[1].text = "--来自「遇见李白」"

    # 内容页
    ic(f'总共{len(ppt_content["pages"])}页')
    for i, page in enumerate(ppt_content['pages']):
        print('生成第%d页:%s' % (i + 1, page['title']))
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # title&content layout
        # 标题
        slide.placeholders[0].text = page['title']
        # 正文
        for sub_content in page['content']:
            print(sub_content)
            # 一级正文
            sub_title = slide.placeholders[1].text_frame.add_paragraph()
            sub_title.text, sub_title.level = sub_content['title'], 1
            # 二级正文
            sub_description = slide.placeholders[1].text_frame.add_paragraph()
            sub_description.text, sub_description.level = sub_content['description'], 2

    _output_file = get_file_path(str(time.time()))
    ppt.save(_output_file)

    return _output_file


def clear_ppt_cache():
    try:
        clear_files_by_timediff(_OUTPUT_DIR, datetime.timedelta(minutes=10).seconds)
    except Exception as e:
        _logger.error(f"clear_ppt_cache error: {e}")


get_scheduler().add_job(clear_ppt_cache, "interval", seconds=10, id="clear_ppt_cache")
